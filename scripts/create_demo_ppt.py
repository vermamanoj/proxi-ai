#!/usr/bin/env python3
"""
Create demo PPT for Proxi hackathon demo.
This PPT represents an existing proposal that Sarah sent to ACME Corporation.
Proxi will add a business case slide to this PPT during the demo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
BRAND_BLUE = RGBColor(0, 82, 147)  # Dark blue
BRAND_CYAN = RGBColor(0, 180, 216)  # Accent cyan
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_slide(prs):
    """Slide 1: Title slide with branding"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_BLUE
    bg.line.fill.background()
    
    # Company name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TechSolutions"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.8))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Software Solutions"
    p.font.size = Pt(28)
    p.font.color.rgb = BRAND_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_proposal_overview(prs):
    """Slide 2: Proposal Overview for ACME Corporation"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_BLUE
    header.line.fill.background()
    
    # Header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "PROPOSAL: ACME Corporation Enterprise Renewal"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Main content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Proposal details
    details = [
        ("Client:", "ACME Corporation"),
        ("Account Tier:", "Enterprise"),
        ("Contract Period:", "June 2026 - May 2029 (3 years)"),
        ("Seats:", "750 users"),
        ("", ""),
        ("Proposed Pricing:", ""),
        ("  • Base Price:", "$2,450/seat/year"),
        ("  • Volume Discount:", "22% (500+ seats)"),
        ("  • Your Price:", "$1,911/seat/year"),
        ("  • Annual Value:", "$1,433,250"),
        ("  • Margin:", "19%"),
    ]
    
    for i, (label, value) in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if label:
            p.text = f"{label} {value}"
        else:
            p.text = ""
        
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        if label and not label.startswith(" "):
            p.font.bold = True
    
    # Date stamp
    date_box = slide.shapes.add_textbox(Inches(10), Inches(6.8), Inches(3), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Prepared: January 15, 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT
    
    return slide

def add_value_proposition(prs):
    """Slide 3: Why TechSolutions"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_BLUE
    header.line.fill.background()
    
    # Header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Why TechSolutions?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Value props
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    props = [
        "✓ 3+ Years Partnership with ACME Corporation",
        "✓ 99.9% Uptime SLA Guarantee",
        "✓ Dedicated Enterprise Support (24/7)",
        "✓ Custom Integration with Your Legacy Systems",
        "✓ Platinum Client Benefits & Priority Roadmap Input",
        "✓ Proven ROI: $4.7M Lifetime Value Partnership",
    ]
    
    for i, prop in enumerate(props):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = prop
        p.font.size = Pt(26)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(18)
    
    return slide

def add_next_steps(prs):
    """Slide 4: Next Steps"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_BLUE
    header.line.fill.background()
    
    # Header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Steps"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Steps content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    steps = [
        "1. Review proposal terms",
        "2. Schedule technical deep-dive (if needed)",
        "3. Finalize contract by March 15, 2026",
        "4. Begin migration planning",
    ]
    
    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(26)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(18)
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(1.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Account Manager"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    
    p = tf.add_paragraph()
    p.text = "Sarah Johnson | sarah.johnson@techsolutions.com | +1-555-0199"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    return slide

# Build the presentation
add_title_slide(prs)
add_proposal_overview(prs)
add_value_proposition(prs)
add_next_steps(prs)

# Save to demo location
output_path = os.path.join(os.path.dirname(__file__), '..', 'demo-apps', 'TechSolutions_ACME_Proposal_2026.pptx')
prs.save(output_path)
print(f"Created: {output_path}")

# Also save to a "Downloads" simulation folder
downloads_path = os.path.join(os.path.dirname(__file__), '..', 'demo-apps', 'downloads')
os.makedirs(downloads_path, exist_ok=True)
downloads_file = os.path.join(downloads_path, 'TechSolutions_ACME_Proposal_2026.pptx')
prs.save(downloads_file)
print(f"Created: {downloads_file}")

print("\nPPT Contents:")
print("  Slide 1: Title - TechSolutions branding")
print("  Slide 2: Proposal Overview - ACME 750 seats, 19% margin, $1,911/seat")
print("  Slide 3: Value Proposition - Why TechSolutions")
print("  Slide 4: Next Steps - Contact Sarah Johnson")
print("\nProxi will add a NEW slide with business case for lower margin approval.")
